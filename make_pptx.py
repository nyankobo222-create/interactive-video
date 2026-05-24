from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── カラーパレット ──────────────────────────────────────────
C_GREEN   = RGBColor(0x2a, 0x98, 0x24)
C_GREEN_L = RGBColor(0xE8, 0xF5, 0xE9)
C_GREEN_D = RGBColor(0x1B, 0x5E, 0x20)
C_DARK    = RGBColor(0x26, 0x32, 0x38)
C_GRAY    = RGBColor(0x54, 0x6E, 0x7A)
C_LIGHT   = RGBColor(0xEC, 0xEF, 0xF1)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLUE    = RGBColor(0x19, 0x76, 0xD2)
C_BLUE_L  = RGBColor(0xE3, 0xF2, 0xFD)
C_ORANGE  = RGBColor(0xF5, 0x7F, 0x17)
C_RED     = RGBColor(0xC6, 0x28, 0x28)
C_PURPLE  = RGBColor(0x6A, 0x1B, 0x9A)
C_NAVY    = RGBColor(0x1A, 0x23, 0x7E)
BG_SLIDE  = RGBColor(0xFA, 0xFD, 0xFA)


def set_bg(slide, rgb):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def box(slide, x, y, w, h, fill_rgb, text="", size=13, bold=False,
        fg=C_WHITE, border_rgb=None, halign=PP_ALIGN.CENTER, rounded=True):
    shape_id = 5 if rounded else 1  # 5=roundRect, 1=rect
    shp = slide.shapes.add_shape(
        shape_id,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    shp.line.fill.background()          # デフォルト枠を消す
    if border_rgb:
        shp.line.color.rgb = border_rgb
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()

    if text:
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left  = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top   = Pt(6)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = halign
            r = p.add_run()
            r.text = line
            r.font.size  = Pt(size)
            r.font.bold  = bold
            r.font.color.rgb = fg
    return shp


def txt(slide, x, y, w, h, text, size=12, bold=False,
        color=C_DARK, halign=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = halign
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = color
    return tb


def arrow(slide, x1, y1, x2, y2, color=C_GRAY, width=2.0):
    cx = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cx.line.color.rgb = color
    cx.line.width = Pt(width)


# ══════════════════════════════════════════════════════════════
# Slide 1: タイトル
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_GREEN_D)

box(sl, 1.2, 1.6, 10.9, 1.6,
    C_GREEN,
    "インタラクティブ動画システム　仕組み・構成 図解",
    size=30, bold=True, fg=C_WHITE, rounded=True)

txt(sl, 1.5, 3.4, 10.3, 0.7,
    "動画ファイル × Webプレーヤー × クラウド配信 の組み合わせで実現する低コスト・高品質なインタラクティブ動画",
    size=14, color=RGBColor(0xC8,0xE6,0xC9), halign=PP_ALIGN.CENTER)

features = [
    ("🎬 動画ファイルを\n差し替えるだけ",   RGBColor(0x2E,0x7D,0x32)),
    ("⚡ 月額ほぼ0円\nで運用できる",        RGBColor(0x1B,0x5E,0x20)),
    ("♾ クライアント数\n無制限で展開",      RGBColor(0x33,0x69,0x1E)),
]
for i, (label, clr) in enumerate(features):
    box(sl, 1.5 + i*3.5, 4.5, 3.0, 1.5, clr, label,
        size=14, fg=C_WHITE, rounded=True)


# ══════════════════════════════════════════════════════════════
# Slide 2: 全体システム構成
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_SLIDE)

txt(sl, 0.4, 0.15, 12.5, 0.55, "全体システム構成",
    size=22, bold=True, color=C_GREEN_D)
txt(sl, 0.4, 0.65, 12.5, 0.4,
    "3つのサービスがそれぞれの役割を担い、連携して動作します",
    size=13, color=C_GRAY)

# GitHub
box(sl, 0.3, 1.3, 2.9, 1.4, RGBColor(0x1A,0x1A,0x2E),
    "GitHub\n📁 コード置き場", size=14, bold=True, fg=C_WHITE)
txt(sl, 0.3, 2.8, 2.9, 0.9,
    "・ソースコード一式\n・変更履歴を管理\n・完全無料", size=10, color=C_GRAY)

# Railway
box(sl, 4.0, 1.3, 3.4, 1.4, C_GREEN,
    "Railway\n🖥 Webサーバー", size=14, bold=True, fg=C_WHITE)
txt(sl, 4.0, 2.8, 3.4, 0.9,
    "・プレーヤーのコードを配信\n・HTMLを返す\n・無料〜$5/月", size=10, color=C_GRAY)

# ユーザー
box(sl, 8.5, 1.3, 2.9, 1.4, C_BLUE,
    "ユーザー端末\n📱 スマホ/PC", size=14, bold=True, fg=C_WHITE)
txt(sl, 8.5, 2.8, 2.9, 0.9,
    "・ブラウザでURLを開く\n・動画を視聴・操作", size=10, color=C_GRAY)

# Cloudflare R2
box(sl, 4.0, 4.2, 3.4, 1.4, C_ORANGE,
    "Cloudflare R2\n🗄 動画ファイル置き場", size=13, bold=True, fg=C_WHITE)
txt(sl, 4.0, 5.7, 3.4, 1.0,
    "・C01〜C14.mp4 を保管\n・世界から高速配信\n・10GB/月まで無料", size=10, color=C_GRAY)

# 矢印とラベル
# GitHub → Railway
arrow(sl, 3.25, 2.0, 3.95, 2.0, C_GREEN, 2.5)
txt(sl, 3.0, 1.6, 1.15, 0.75, "push\n→ 自動デプロイ",
    size=9, color=C_GREEN, halign=PP_ALIGN.CENTER)

# Railway ↔ ユーザー
arrow(sl, 7.45, 1.9, 8.45, 1.9, C_GREEN, 2.5)
arrow(sl, 8.45, 2.15, 7.45, 2.15, C_GRAY, 1.5)
txt(sl, 7.25, 1.5, 1.4, 0.45, "コード配信 →", size=9, color=C_GREEN)
txt(sl, 7.25, 1.95, 1.4, 0.45, "← アクセス",  size=9, color=C_GRAY)

# Railway ↕ R2
arrow(sl, 5.7, 2.75, 5.7, 4.15, C_ORANGE, 2.0)
txt(sl, 5.85, 3.3, 1.8, 0.5, "動画URLを\n参照・取得", size=9, color=C_ORANGE)

# ユーザー → R2
arrow(sl, 9.9, 2.75, 7.5, 4.2, C_ORANGE, 1.5)
txt(sl, 8.5, 3.65, 1.8, 0.6, "MP4ファイルを\n直接ダウンロード",
    size=9, color=C_ORANGE, halign=PP_ALIGN.CENTER)

# ポイントボックス
box(sl, 0.3, 5.15, 3.3, 2.1, C_GREEN_L,
    "✅ ポイント\n\nコードと動画は\n完全に別管理\n\n動画だけ差し替えれば\n内容を更新できる",
    size=11, fg=C_GREEN_D, border_rgb=C_GREEN, halign=PP_ALIGN.LEFT, rounded=True)

box(sl, 11.5, 1.3, 1.5, 5.95, C_GREEN_L,
    "コード\n管理\n\n↕\n\n動画\n管理\n\n↕\n\n分離",
    size=10, fg=C_GREEN_D, border_rgb=C_GREEN, halign=PP_ALIGN.CENTER, rounded=False)


# ══════════════════════════════════════════════════════════════
# Slide 3: UIレイヤーの仕組み
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_SLIDE)

txt(sl, 0.4, 0.15, 12.5, 0.55,
    'UIレイヤーの仕組み ── 「付箋を貼る」イメージ',
    size=22, bold=True, color=C_GREEN_D)
txt(sl, 0.4, 0.65, 12.5, 0.4,
    "動画とボタンは別々のHTMLで、CSSによって重ね合わせているだけです。動画が変わってもボタンUIは影響を受けません。",
    size=13, color=C_GRAY)

# ── 左側: 重なりイメージ（ずらして配置） ──
layers_data = [
    (0.4,  5.2, 4.8, 1.0, C_DARK,   "【外枠】 <div class='player'>",
     "position: relative\n（すべてのレイヤーの「基準点」）"),
    (0.7,  4.1, 4.5, 1.0, C_NAVY,   "【1層目】 <video>",
     "動画レイヤー（最背面）\nMP4をただ再生するだけ"),
    (1.0,  3.0, 4.5, 1.0, C_GREEN,  "【2層目】 <BottomBar>",
     "ボトムナビバー\nposition: absolute / bottom: 0"),
    (1.3,  1.9, 4.5, 1.0, C_PURPLE, "【3層目】 <BranchMenu>",
     "ブランチ選択UI（最前面）\nposition: absolute / inset: 0"),
]
for x, y, w, h, clr, title, sub in layers_data:
    box(sl, x, y, w, h, clr, title, size=12, bold=True, fg=C_WHITE, rounded=False)
    txt(sl, x + 4.85, y+0.05, 2.8, 0.9, sub, size=10, color=C_GRAY)

txt(sl, 0.1, 1.7,  0.6, 0.4, "▲前面", size=9, color=C_PURPLE, halign=PP_ALIGN.CENTER)
txt(sl, 0.1, 5.5,  0.6, 0.4, "▼背面", size=9, color=C_GRAY,   halign=PP_ALIGN.CENTER)
arrow(sl, 0.35, 5.45, 0.35, 1.9, C_GRAY, 1.0)

# ── 右側: CSS説明 ──
box(sl, 8.0, 1.4, 5.0, 1.55, C_GREEN_L,
    "【たとえ話】付箋を貼るイメージ\n\nスライド（動画）はそのまま流れている\n↓\nその上に付箋（ボタン）を貼っているだけ",
    size=11, fg=C_GREEN_D, border_rgb=C_GREEN, halign=PP_ALIGN.LEFT, rounded=True)

box(sl, 8.0, 3.15, 5.0, 1.55, C_NAVY,
    "position: absolute とは？\n\n「親の箱を基準に、自由な位置に\n浮かせる」というCSS命令",
    size=11, fg=C_WHITE, halign=PP_ALIGN.LEFT, rounded=True)

box(sl, 8.0, 4.9, 5.0, 1.55, C_PURPLE,
    "inset: 0 とは？\n\n「上下左右すべて0」= 全面を覆う。\nこれで動画の上を完全に覆える。",
    size=11, fg=C_WHITE, halign=PP_ALIGN.LEFT, rounded=True)

box(sl, 0.4, 6.55, 12.5, 0.75, C_GREEN,
    "✅ 動画とUIは完全独立！　→　動画ファイルをどう変えてもボタンのデザインや動作は一切影響を受けない",
    size=12, bold=True, fg=C_WHITE, rounded=True)


# ══════════════════════════════════════════════════════════════
# Slide 4: 動作フロー
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_SLIDE)

txt(sl, 0.4, 0.15, 12.5, 0.55, "プレーヤーの動作フロー（状態遷移図）",
    size=22, bold=True, color=C_GREEN_D)
txt(sl, 0.4, 0.65, 12.5, 0.4,
    "プレーヤーは常に4つの「状態」のどれかにあり、操作・時間経過・動画終了で次の状態へ移行します",
    size=13, color=C_GRAY)

# 状態ボックス
states = [
    (0.4,  1.5, C_BLUE,   "① intro\nイントロ再生中",     "C01が自動再生される"),
    (4.15, 1.5, C_GREEN,  "② branch_select\n選択待ち",   "3ボタンが大きく表示される"),
    (7.9,  1.5, C_ORANGE, "③ branch_playing\n再生中",    "選択ブランチが順番に流れる"),
    (4.15, 4.5, C_PURPLE, "④ end_menu\nエンドメニュー",  "C14後、全ボタン再表示"),
]
for x, y, clr, title, sub in states:
    box(sl, x, y, 3.3, 1.5, clr, title, size=14, bold=True, fg=C_WHITE)
    txt(sl, x, y+1.6, 3.3, 0.5, sub, size=10, color=C_GRAY, halign=PP_ALIGN.CENTER)

# 矢印
arrow(sl, 3.75, 2.25, 4.1,  2.25, C_BLUE,   2.5)
arrow(sl, 7.5,  2.25, 7.85, 2.25, C_GREEN,  2.5)
arrow(sl, 9.55, 3.05, 7.0,  4.45, C_ORANGE, 2.5)
arrow(sl, 5.85, 4.45, 8.7,  3.05, C_PURPLE, 2.0)
arrow(sl, 4.15, 5.25, 1.85, 3.05, C_PURPLE, 2.0)

# 矢印ラベル
txt(sl, 3.3,  1.85, 1.0, 0.7,  "15秒\n経過→",      size=9, color=C_BLUE,   halign=PP_ALIGN.CENTER)
txt(sl, 7.15, 1.85, 0.9, 0.7,  "ボタン\nタップ→",   size=9, color=C_GREEN,  halign=PP_ALIGN.CENTER)
txt(sl, 8.05, 3.5,  1.5, 0.7,  "全チャプター\n終了", size=9, color=C_ORANGE, halign=PP_ALIGN.CENTER)
txt(sl, 6.8,  3.9,  1.5, 0.7,  "別ブランチ\n選択",   size=9, color=C_PURPLE, halign=PP_ALIGN.CENTER)
txt(sl, 2.1,  4.25, 1.8, 0.7,  "「トップへ\n戻る」", size=9, color=C_PURPLE, halign=PP_ALIGN.CENTER)

# FKENの具体例
box(sl, 0.4, 5.85, 12.5, 1.4, C_GREEN_L,
    "【FKENの場合の具体的な流れ】\n"
    "C01再生（15秒）→ 停止 → ボタン3つ表示 → 「01 未経験からプロへ」をタップ\n"
    "→ C02 → C03 → C04 → C05 と自動再生 → C14（エンドメニュー）→ 他ブランチを選ぶかトップへ",
    size=11, fg=C_GREEN_D, border_rgb=C_GREEN, halign=PP_ALIGN.LEFT, rounded=True)


# ══════════════════════════════════════════════════════════════
# Slide 5: 設定ファイルとプレーヤーの関係
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_SLIDE)

txt(sl, 0.4, 0.15, 12.5, 0.55, "設定ファイルとプレーヤーの関係",
    size=22, bold=True, color=C_GREEN_D)
txt(sl, 0.4, 0.65, 12.5, 0.4,
    "「何をどう動かすか（設定）」と「動かす仕組み（エンジン）」が分離しているため、設定だけ変えれば別会社に展開できます",
    size=13, color=C_GRAY)

# 左: 設定ファイル
box(sl, 0.3, 1.35, 4.5, 5.45, C_DARK,
    "fken.js（設定ファイル）\n\n"
    "company.name\n  = 株式会社FKEN\n\n"
    "theme.primary\n  = #2a9824（緑）\n\n"
    "chapters.C01.url\n  = https://r2.dev/.../C01.mp4\n\n"
    "flow.branches[0]\n  = label: 01 未経験からプロへ\n  chapters: C02,C03,C04,C05\n\n"
    "flow.intro.pauseAt\n  = 15（秒）",
    size=11, fg=C_WHITE, halign=PP_ALIGN.LEFT, rounded=True)

# 矢印と意味
pairs = [
    (2.05, "会社名をエンドメニューに表示する"),
    (2.85, "全ボタンをこの色で塗る"),
    (3.65, "C01として再生する動画ファイル"),
    (4.5,  "ボタンラベルと再生順序を決める"),
    (5.35, "何秒でC01を止めてメニューを出すか"),
]
for y, label in pairs:
    arrow(sl, 4.85, y, 5.35, y, C_GREEN, 2.0)
    box(sl, 5.4, y-0.28, 4.3, 0.6, C_GREEN_L, label,
        size=11, fg=C_GREEN_D, halign=PP_ALIGN.LEFT, rounded=True)

# 右: プレーヤー本体
box(sl, 9.85, 1.35, 3.1, 5.45, C_GREEN,
    "InteractivePlayer.jsx\n（プレーヤー本体）\n\n設定通りに動く\n汎用エンジン\n\n"
    "✓ 会社名を表示\n✓ その色でボタン\n✓ そのURLの動画\n✓ そのラベルで分岐\n✓ その秒数で停止",
    size=12, fg=C_WHITE, halign=PP_ALIGN.LEFT, rounded=True)

box(sl, 0.3, 7.0, 12.7, 0.35, C_GREEN,
    "💡 新クライアントが来たら fken.js をコピーして会社名・色・URL・ラベルを書き換えるだけ。エンジン（プレーヤー）は一切触らない。",
    size=11, bold=True, fg=C_WHITE, rounded=True)


# ══════════════════════════════════════════════════════════════
# Slide 6: 複数クライアント展開
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_SLIDE)

txt(sl, 0.4, 0.15, 12.5, 0.55, "複数クライアントへの展開イメージ",
    size=22, bold=True, color=C_GREEN_D)
txt(sl, 0.4, 0.65, 12.5, 0.4,
    "1つのシステムで何社でも対応。クライアントが増えても追加コストはほぼゼロ。",
    size=13, color=C_GRAY)

# 中央: エンジン
box(sl, 4.65, 3.1, 4.0, 1.6, C_GREEN,
    "InteractivePlayer\n（プレーヤー本体）\n1つだけ作る・共通で使う",
    size=13, bold=True, fg=C_WHITE)

# 周囲: 設定ファイル × 5社
clients = [
    (0.3,  0.85, C_DARK,                    "fken.js\n株式会社FKEN\n緑・電話番号A"),
    (4.65, 0.85, RGBColor(0x0D,0x47,0xA1),  "abc.js\n株式会社ABC\n青・電話番号B"),
    (9.1,  0.85, C_PURPLE,                  "xyz.js\n株式会社XYZ\n紫・電話番号C"),
    (0.3,  5.3,  C_ORANGE,                  "def.js\n株式会社DEF\nオレンジ・電話番号D"),
    (9.1,  5.3,  RGBColor(0x2E,0x7D,0x32),  "ghi.js\n株式会社GHI\n深緑・電話番号E"),
]
targets = [
    (3.3, 1.6, 4.65, 3.4),
    (6.65,1.6, 6.65, 3.1),
    (9.1, 1.6, 8.65, 3.4),
    (3.3, 5.8, 4.65, 4.7),
    (9.1, 5.8, 8.65, 4.7),
]
for (x,y,clr,lbl), (x1,y1,x2,y2) in zip(clients, targets):
    box(sl, x, y, 3.4, 1.4, clr, lbl, size=11, fg=C_WHITE)
    arrow(sl, x1, y1, x2, y2, C_GRAY, 1.5)

box(sl, 0.3, 6.9, 12.7, 0.45, C_GREEN,
    "各クライアントのURL例：　yourapp.com/fken　yourapp.com/abc　yourapp.com/xyz　→ 設定ファイル1枚で何社でも展開可能",
    size=11, bold=True, fg=C_WHITE, rounded=True)


# ══════════════════════════════════════════════════════════════
# Slide 7: コスト比較
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_SLIDE)

txt(sl, 0.4, 0.15, 12.5, 0.55, "コスト比較：ミル vs 自作システム",
    size=22, bold=True, color=C_GREEN_D)
txt(sl, 0.4, 0.65, 12.5, 0.4,
    "クライアント数が増えるほど自作の優位性が拡大します",
    size=13, color=C_GRAY)

# ヘッダー行
box(sl, 0.4,  1.3, 4.4, 0.65, C_DARK,  "比較項目",             size=13, bold=True, fg=C_WHITE, rounded=False)
box(sl, 4.85, 1.3, 3.7, 0.65, C_RED,   "ミル（既存サービス）",  size=13, bold=True, fg=C_WHITE, rounded=False)
box(sl, 8.6,  1.3, 4.5, 0.65, C_GREEN, "自作システム（今回）",  size=13, bold=True, fg=C_WHITE, rounded=False)

rows = [
    ("初期費用",               "プラン契約料（高め）",          "開発工数のみ（構築済み）",       False),
    ("月額（1社）",             "数万円 / 月",                  "ほぼ 0円 / 月",                 True),
    ("月額（5社）",             "数万円 × 5",                   "変わらず ほぼ 0円",              True),
    ("動画ストレージ",          "プランに含む（上限あり）",       "R2：10GB/月まで無料",            False),
    ("デザイン・機能",          "制限あり（テンプレ内）",         "完全自由",                      False),
    ("データ所有権",            "サービス側",                    "自分（GitHub管理）",             False),
    ("サービス終了リスク",       "あり",                          "なし（自前サーバー）",           False),
]
for i, (item, mill, ours, highlight) in enumerate(rows):
    y = 2.05 + i * 0.65
    bg_i = C_LIGHT
    bg_m = RGBColor(0xFF,0xEB,0xEE) if highlight else C_WHITE
    bg_o = RGBColor(0xE8,0xF5,0xE9) if highlight else C_WHITE
    fc_m = C_RED   if highlight else C_DARK
    fc_o = C_GREEN if highlight else C_DARK
    box(sl, 0.4,  y, 4.4, 0.62, bg_i, item, size=12, fg=C_DARK,  rounded=False, border_rgb=C_LIGHT)
    box(sl, 4.85, y, 3.7, 0.62, bg_m, mill, size=11, fg=fc_m, bold=highlight, rounded=False, border_rgb=C_LIGHT)
    box(sl, 8.6,  y, 4.5, 0.62, bg_o, ours, size=11, fg=fc_o, bold=highlight, rounded=False, border_rgb=C_LIGHT)

box(sl, 0.4, 6.65, 12.5, 0.7, C_GREEN,
    "💡 5社での月額差は数十万円以上になる可能性も。クライアントが増えるほど回収が早くなります。",
    size=12, bold=True, fg=C_WHITE, rounded=True)


# ══════════════════════════════════════════════════════════════
# 保存
# ══════════════════════════════════════════════════════════════
out = r"D:\interactive-video\インタラクティブ動画システム_図解.pptx"
prs.save(out)
print(f"保存完了: {out}")
